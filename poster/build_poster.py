#!/usr/bin/env python3
"""Build the FM-AVSR one-page poster from checked-in assets."""

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
POSTER_DIR = ROOT / "poster"
ASSET_DIR = POSTER_DIR / "assets"
TEMPLATE = POSTER_DIR / "Poster Template.pptx"
OUT = POSTER_DIR / "fm_avsr_poster.pptx"


NAVY = RGBColor(18, 35, 61)
TEAL = RGBColor(14, 116, 144)
GREEN = RGBColor(22, 163, 74)
ORANGE = RGBColor(234, 88, 12)
RED = RGBColor(185, 28, 28)
SLATE = RGBColor(51, 65, 85)
MUTED = RGBColor(100, 116, 139)
LIGHT = RGBColor(248, 250, 252)
BORDER = RGBColor(203, 213, 225)
WHITE = RGBColor(255, 255, 255)

# Font size scale for a 33×43 inch poster
SZ_BODY    = 40   # regular body text
SZ_SMALL   = 32   # captions, table cells, fine print
SZ_SUB     = 52   # panel section sub-labels, emphasis
SZ_PANEL   = 60   # panel titles (small-caps style)
SZ_SECTION = 80   # section / slide title
SZ_HERO    = 100  # hero numbers / big claims


def inch(value: float):
    return Inches(value)


def clear_slide(slide):
    sp_tree = slide.shapes._spTree
    for element in list(sp_tree):
        tag = element.tag.rsplit("}", 1)[-1]
        if tag in {"nvGrpSpPr", "grpSpPr"}:
            continue
        sp_tree.remove(element)


def clear_body_only(slide, slide_height_inches: float):
    """Keep template header/footer artwork; remove editable placeholder body."""
    sp_tree = slide.shapes._spTree
    keep = []
    for shape in slide.shapes:
        y = shape.top / 914400
        bottom = (shape.top + shape.height) / 914400
        text = shape.text.strip() if hasattr(shape, "text") else ""
        is_template_header = y < 2.35
        is_template_footer = bottom > slide_height_inches - 2.05
        is_course_footer = "Deep Learning 2026 Spring" in text
        keep.append(is_template_header or is_template_footer or is_course_footer)

    for element, should_keep in zip(list(sp_tree)[2:], keep):
        if should_keep:
            continue
        sp_tree.remove(element)


def set_text(tf, text, size=SZ_BODY, bold=False, color=SLATE, align=None):
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Arial"
    if align is not None:
        p.alignment = align


def add_text(slide, text, x, y, w, h, size=SZ_BODY, bold=False, color=SLATE, align=None):
    box = slide.shapes.add_textbox(inch(x), inch(y), inch(w), inch(h))
    box.text_frame.word_wrap = True
    box.text_frame.margin_left = inch(0.08)
    box.text_frame.margin_right = inch(0.08)
    box.text_frame.margin_top = inch(0.04)
    box.text_frame.margin_bottom = inch(0.04)
    set_text(box.text_frame, text, size, bold, color, align)
    return box


def add_panel(slide, title, x, y, w, h, accent=TEAL, fill=WHITE):
    panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, inch(x), inch(y), inch(w), inch(h)
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = fill
    panel.line.color.rgb = BORDER
    panel.line.width = Pt(1.0)
    panel.adjustments[0] = 0.035
    add_text(slide, title, x + 0.22, y + 0.18, w - 0.44, 0.55, SZ_PANEL, True, accent)
    return panel


def add_tag(slide, text, x, y, w, color=TEAL):
    tag = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, inch(x), inch(y), inch(w), inch(0.72)
    )
    tag.fill.solid()
    tag.fill.fore_color.rgb = color
    tag.line.fill.background()
    tag.adjustments[0] = 0.14
    add_text(slide, text, x + 0.12, y + 0.12, w - 0.24, 0.48, SZ_BODY, True, WHITE, PP_ALIGN.CENTER)


def add_claim_chip(slide, label, body, x, y, w, accent):
    chip = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, inch(x), inch(y), inch(w), inch(0.62)
    )
    chip.fill.solid()
    chip.fill.fore_color.rgb = RGBColor(248, 250, 252)
    chip.line.color.rgb = RGBColor(186, 230, 253)
    chip.adjustments[0] = 0.08
    marker = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, inch(x), inch(y), inch(0.08), inch(0.62)
    )
    marker.fill.solid()
    marker.fill.fore_color.rgb = accent
    marker.line.fill.background()
    add_text(slide, label, x + 0.18, y + 0.08, 1.35, 0.18, 8, True, accent)
    add_text(slide, body, x + 1.48, y + 0.08, w - 1.62, 0.32, 9, True, NAVY)


def add_artifact_card(slide, title, body, x, y, w, h, accent=TEAL):
    card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, inch(x), inch(y), inch(w), inch(h)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(248, 250, 252)
    card.line.color.rgb = RGBColor(226, 232, 240)
    card.adjustments[0] = 0.06
    add_text(slide, title, x + 0.14, y + 0.14, w - 0.28, 0.38, SZ_BODY, True, accent)
    add_text(slide, body,  x + 0.14, y + 0.60, w - 0.28, h - 0.70, SZ_SMALL, False, NAVY)


def add_intro_card(slide, label, body, x, y, w, h, accent=TEAL):
    card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, inch(x), inch(y), inch(w), inch(h)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(248, 250, 252)
    card.line.color.rgb = RGBColor(226, 232, 240)
    card.adjustments[0] = 0.06
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, inch(x), inch(y), inch(0.12), inch(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    add_text(slide, label, x + 0.25, y + 0.12, w - 0.38, 0.38, SZ_BODY, True, accent)
    add_text(slide, body,  x + 0.25, y + 0.60, w - 0.38, h - 0.68, SZ_SMALL, False, NAVY)


def add_metric_row(slide, label, values, x, y, widths, size=SZ_SMALL):
    add_text(slide, label, x, y, widths[0], 0.24, size, True, NAVY)
    cursor = x + widths[0]
    for value, width in zip(values, widths[1:]):
        add_text(slide, value, cursor, y, width, 0.24, size, False, SLATE, PP_ALIGN.RIGHT)
        cursor += width


def add_mini_table(slide, headers, rows, x, y, w, row_h=0.55, size=SZ_SMALL):
    widths = [w * 0.49, w * 0.17, w * 0.17, w * 0.17]
    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, inch(x), inch(y), inch(w), inch(row_h))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(236, 253, 252)
    bg.line.color.rgb = RGBColor(186, 230, 253)
    add_metric_row(slide, headers[0], headers[1:], x + 0.08, y + 0.07, widths, size)
    for idx, row in enumerate(rows):
        yy = y + row_h * (idx + 1)
        fill = RGBColor(248, 250, 252) if idx % 2 == 0 else WHITE
        rect = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, inch(x), inch(yy), inch(w), inch(row_h))
        rect.fill.solid()
        rect.fill.fore_color.rgb = fill
        rect.line.color.rgb = RGBColor(226, 232, 240)
        add_metric_row(slide, row[0], row[1:], x + 0.08, yy + 0.07, widths, size)


def add_progression_chart(slide, x, y, w, h):
    labels = ["10k", "30k", "timbre", "prompt+res", "59k"]
    values = [0.507, 0.531, 0.563, 0.573, 0.582]
    max_v = 0.60
    bar_gap = 0.18
    bar_w = (w - bar_gap * (len(values) - 1)) / len(values)
    for i, (label, value) in enumerate(zip(labels, values)):
        bx = x + i * (bar_w + bar_gap)
        bh = h * value / max_v
        by = y + h - bh
        color = GREEN if i == len(values) - 1 else (TEAL if i >= 2 else RGBColor(59, 130, 246))
        bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, inch(bx), inch(by), inch(bar_w), inch(bh))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        add_text(slide, f"{value:.3f}", bx - 0.03, by - 0.55, bar_w + 0.06, 0.45, SZ_SMALL, True, NAVY, PP_ALIGN.CENTER)
        add_text(slide, label, bx - 0.05, y + h + 0.1, bar_w + 0.1, 0.45, SZ_SMALL, False, MUTED, PP_ALIGN.CENTER)
    axis = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, inch(x), inch(y + h), inch(w), inch(0.02))
    axis.fill.solid()
    axis.fill.fore_color.rgb = BORDER
    axis.line.fill.background()


def add_image(slide, path, x, y, w, h):
    if path.exists():
        return slide.shapes.add_picture(str(path), inch(x), inch(y), inch(w), inch(h))
    missing = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, inch(x), inch(y), inch(w), inch(h)
    )
    missing.fill.solid()
    missing.fill.fore_color.rgb = RGBColor(226, 232, 240)
    missing.line.color.rgb = BORDER
    add_text(slide, f"Missing asset:\n{path.name}", x + 0.2, y + 0.2, w - 0.4, h - 0.4, 12, True, MUTED, PP_ALIGN.CENTER)
    return missing


def add_fit_image(slide, path, x, y, w, h, mode="contain"):
    """Place an image without distorting its aspect ratio."""
    if not path.exists():
        return add_image(slide, path, x, y, w, h)
    with Image.open(path) as img:
        ratio = img.width / img.height
    box_ratio = w / h
    if (mode == "cover" and ratio < box_ratio) or (mode == "contain" and ratio > box_ratio):
        draw_w = w
        draw_h = w / ratio
    else:
        draw_h = h
        draw_w = h * ratio
    draw_x = x + (w - draw_w) / 2
    draw_y = y + (h - draw_h) / 2
    return slide.shapes.add_picture(str(path), inch(draw_x), inch(draw_y), inch(draw_w), inch(draw_h))


def add_frame_strip(slide, title, paths, x, y, w, h, accent=TEAL):
    add_text(slide, title, x, y, w, 0.45, SZ_BODY, True, accent, PP_ALIGN.CENTER)
    gap = 0.08
    frame_w = (w - gap * (len(paths) - 1)) / len(paths)
    for idx, path in enumerate(paths):
        fx = x + idx * (frame_w + gap)
        bg = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, inch(fx), inch(y + 0.34), inch(frame_w), inch(h - 0.34)
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = WHITE
        bg.line.color.rgb = RGBColor(226, 232, 240)
        bg.adjustments[0] = 0.04
        add_fit_image(slide, path, fx + 0.03, y + 0.37, frame_w - 0.06, h - 0.4)


def add_section(slide, title, x, y, w, h):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, inch(x), inch(y), inch(w), inch(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = BORDER
    shape.line.width = Pt(1.2)
    shape.adjustments[0] = 0.06
    add_text(slide, title, x + 0.25, y + 0.22, w - 0.5, 0.7, SZ_SECTION, True, NAVY)
    rule = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, inch(x + 0.25), inch(y + 1.05), inch(w - 0.5), inch(0.05)
    )
    rule.fill.solid()
    rule.fill.fore_color.rgb = TEAL
    rule.line.fill.background()
    return shape


def add_bullets(slide, items, x, y, w, h, size=SZ_BODY, color=SLATE):
    box = slide.shapes.add_textbox(inch(x), inch(y), inch(w), inch(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = inch(0.05)
    tf.margin_right = inch(0.05)
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(size)
        p.font.name = "Arial"
        p.font.color.rgb = color
        p.space_after = Pt(5)
        p.line_spacing = 1.05
    return box


def add_metric(slide, label, value, note, x, y, w, color):
    card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, inch(x), inch(y), inch(w), inch(1.35)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(241, 245, 249)
    card.line.color.rgb = BORDER
    card.adjustments[0] = 0.08
    add_text(slide, value, x + 0.15, y + 0.08, w - 0.3, 0.48, 25, True, color, PP_ALIGN.CENTER)
    add_text(slide, label, x + 0.15, y + 0.58, w - 0.3, 0.32, 11, True, NAVY, PP_ALIGN.CENTER)
    add_text(slide, note, x + 0.15, y + 0.91, w - 0.3, 0.32, 10, False, MUTED, PP_ALIGN.CENTER)


def add_captioned_image(slide, path, caption, x, y, w, h):
    if path.exists():
        slide.shapes.add_picture(str(path), inch(x), inch(y), inch(w), inch(h))
    else:
        missing = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE, inch(x), inch(y), inch(w), inch(h)
        )
        missing.fill.solid()
        missing.fill.fore_color.rgb = RGBColor(226, 232, 240)
        missing.line.color.rgb = BORDER
    add_text(slide, caption, x, y + h + 0.1, w, 0.55, SZ_BODY, True, MUTED, PP_ALIGN.CENTER)


def add_bar_chart(slide, x, y, w, h):
    labels = ["10k", "30k", "+timbre", "+residual", "59k final"]
    values = [0.5073, 0.5308, 0.5633, 0.5726, 0.5818]
    max_v = 0.60
    bar_w = w / len(values) * 0.62
    gap = w / len(values) * 0.38
    for i, (label, value) in enumerate(zip(labels, values)):
        bx = x + i * (bar_w + gap) + gap * 0.5
        bh = h * value / max_v
        by = y + h - bh
        bar = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE, inch(bx), inch(by), inch(bar_w), inch(bh)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = GREEN if i == len(values) - 1 else TEAL
        bar.line.fill.background()
        add_text(slide, f"{value:.3f}", bx - 0.05, by - 0.65, bar_w + 0.1, 0.5, SZ_SMALL, True, NAVY, PP_ALIGN.CENTER)
        add_text(slide, label, bx - 0.15, y + h + 0.12, bar_w + 0.3, 0.55, SZ_SMALL, False, MUTED, PP_ALIGN.CENTER)
    axis = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, inch(x), inch(y + h), inch(w), inch(0.025)
    )
    axis.fill.solid()
    axis.fill.fore_color.rgb = BORDER
    axis.line.fill.background()


def add_pipeline(slide, x, y, w, h):
    steps = [
        ("silent video", TEAL),
        ("AVSR visual\n+ weak text", RGBColor(37, 99, 235)),
        ("timbre ref\n3.04 s", ORANGE),
        ("residual\nMimi recon", GREEN),
        ("speech video", NAVY),
    ]
    box_w = w / 5.0 - 0.12
    for i, (label, color) in enumerate(steps):
        bx = x + i * (box_w + 0.15)
        node = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, inch(bx), inch(y), inch(box_w), inch(h)
        )
        node.fill.solid()
        node.fill.fore_color.rgb = color
        node.line.fill.background()
        node.adjustments[0] = 0.12
        add_text(slide, label, bx + 0.05, y + 0.18, box_w - 0.1, h - 0.2, 15, True, WHITE, PP_ALIGN.CENTER)
        if i < len(steps) - 1:
            add_text(slide, ">", bx + box_w - 0.01, y + 0.33, 0.22, 0.32, 16, True, MUTED, PP_ALIGN.CENTER)


def add_condition_table(slide, x, y, w, h):
    rows = [
        ("lip_avsr.npy", "Auto-AVSR lip crop latent"),
        ("avsr_enc_lipavsr.npy", "AVSR visual encoder state"),
        ("avsr_text_lipavsr.txt", "noisy text, weak semantic aid"),
        ("SmolLM2 hidden", "aligned text-token hidden states"),
        ("speaker/timbre", "speaker_emb + Mimi prompt stats"),
        ("target", "Mimi audio latent to reconstruct"),
    ]
    row_h = h / len(rows)
    for i, (name, desc) in enumerate(rows):
        ry = y + i * row_h
        bg = RGBColor(248, 250, 252) if i % 2 == 0 else WHITE
        rect = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, inch(x), inch(ry), inch(w), inch(row_h))
        rect.fill.solid()
        rect.fill.fore_color.rgb = bg
        rect.line.color.rgb = RGBColor(226, 232, 240)
        add_text(slide, name, x + 0.12, ry + 0.08, 2.3, row_h - 0.12, SZ_SMALL, True, NAVY)
        add_text(slide, desc, x + 2.45, ry + 0.08, w - 2.55, row_h - 0.12, SZ_SMALL, False, SLATE)


def build():
    prs = Presentation(str(TEMPLATE))
    slide = prs.slides[0]

    sw = prs.slide_width / 914400
    sh = prs.slide_height / 914400
    clear_body_only(slide, sh)

    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = LIGHT

    # Body background. Template header/footer artwork remains untouched.
    body_top = 2.45
    body_bottom = 41.18
    body_h = body_bottom - body_top
    bg_rect = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, inch(0), inch(body_top), inch(sw), inch(body_h)
    )
    bg_rect.fill.solid()
    bg_rect.fill.fore_color.rgb = RGBColor(247, 251, 252)
    bg_rect.line.fill.background()

    # Title block.
    add_text(slide, "Audio-Centric Visual Speech Reconstruction", 1.05, 2.78, 31.0, 1.2, SZ_SECTION, True, NAVY, PP_ALIGN.CENTER)
    add_text(
        slide,
        "Generate speech by reconstructing Mimi audio latents, not by cascading lip → text → TTS",
        2.1, 4.1, 28.9, 0.8, SZ_SUB, True, TEAL, PP_ALIGN.CENTER,
    )
    # Central architecture
    top_y = 5.2
    add_panel(slide, "1  CORE ARCHITECTURE", 1.05, top_y, 22.0, 13.5)
    add_fit_image(slide, ASSET_DIR / "report_system_architecture.png", 1.45, top_y + 0.75, 21.2, 10.0)
    add_text(slide, "Conditions: lip/video  |  Weak text hidden states  |  Speaker embedding  |  3.04s Mimi audio prompt  |  Prompt timbre stats",
             1.45, top_y + 10.95, 21.2, 0.65, SZ_SMALL, True, TEAL, PP_ALIGN.CENTER)
    add_text(slide, "Target: normalized Mimi latents  →  waveform via Mimi decoder  (offline, 24 kHz mono)",
             1.45, top_y + 11.7, 21.2, 0.65, SZ_SMALL, False, SLATE, PP_ALIGN.CENTER)
    add_text(slide, "Key: text is NOT the main generation path — lip and audio-prompt latents carry most of the signal",
             1.45, top_y + 12.45, 21.2, 0.75, SZ_SMALL, True, ORANGE, PP_ALIGN.CENTER)

    add_panel(slide, "2  PROJECT OVERVIEW", 23.45, top_y, 8.6, 7.2)
    add_intro_card(slide, "main claim",
        "Do not turn visual speech generation into a strict lip → text → TTS cascade. Text is a weak semantic/alignment cue.",
        23.78, top_y + 0.88, 7.9, 1.8, TEAL)
    add_intro_card(slide, "route",
        "Lip motion, weak LM text, speaker identity, and reference-timbre prompt are fused to reconstruct continuous Mimi audio latents.",
        23.78, top_y + 2.88, 7.9, 1.92, ORANGE)
    add_intro_card(slide, "evidence",
        "Trained on 59,144 LRS3 lip-AVSR clips; final checkpoint reaches corr 0.5818 / MSE 0.6676 / MAE 0.6018.",
        23.78, top_y + 5.0, 7.9, 1.45, GREEN)

    add_panel(slide, "3  DATA + TRAINING SIGNALS", 23.45, top_y + 7.75, 8.6, 5.5)
    add_text(slide, "59,144 lip-AVSR clips  |  fixed 1,000 val", 23.82, top_y + 8.55, 7.9, 0.55, SZ_BODY, True, NAVY)
    add_condition_table(slide, 23.78, top_y + 9.25, 7.85, 3.1)
    add_tag(slide, "lip crops + weak text + timbre prompt", 24.25, top_y + 12.55, 6.9, TEAL)

    # Middle row: text ablation / residual / experimental
    row_y = 19.4
    gap = 0.35
    col_w = (31.0 - 2 * gap) / 3
    x1 = 1.05
    x2 = x1 + col_w + gap
    x3 = x2 + col_w + gap

    add_panel(slide, "4  TEXT ABLATION", x1, row_y, col_w, 7.95)
    add_text(slide, "GT transcript is not the core bottleneck.", x1 + 0.35, row_y + 0.85, col_w - 0.7, 0.65, SZ_SUB, True, TEAL)
    rows = [
        ("Text source", "Corr", "MSE", "MAE"),
        ("GT text + timestamps",   "0.5818", "0.6676", "0.6018"),
        ("AVSR text + timestamps", "0.5783", "0.6716", "0.6037"),
        ("AVSR text + uniform",    "0.2868", "1.684",  "1.024"),
        ("V5 text + uniform",      "0.2861", "1.685",  "1.024"),
    ]
    add_mini_table(slide, rows[0], rows[1:], x1 + 0.4, row_y + 1.75, col_w - 0.8, row_h=0.55, size=SZ_SMALL)
    add_text(
        slide,
        "AVSR vs GT: Δcorr = 0.0035 only.\nV5 (WER 29%) vs AVSR (WER 20%): Δcorr = 0.0007.\nAlignment quality (timestamps vs uniform) = Δ0.29.",
        x1 + 0.32, row_y + 4.75, col_w - 0.64, 1.5, SZ_BODY, False, SLATE,
    )
    add_text(slide, "alignment > accuracy", x1 + 0.55, row_y + 6.3, col_w - 1.1, 0.95, SZ_SECTION, True, GREEN, PP_ALIGN.CENTER)

    add_panel(slide, "5  RESIDUAL RECON", x2, row_y, col_w, 7.95)
    add_fit_image(slide, ASSET_DIR / "report_recon_vs_sampling.png", x2 + 0.35, row_y + 0.85, col_w - 0.7, 3.6)
    add_text(
        slide,
        "FM sampling (NFE=10) stalls at corr ≈ 0.35. Fix: set x̃=0, τ=1 — predict endpoint directly.",
        x2 + 0.38, row_y + 4.6, col_w - 0.76, 0.9, SZ_BODY, True, NAVY,
    )
    add_text(
        slide,
        "Frozen base reconstructor predicts latent. Trainable residual head corrects errors. Doubles corr vs naive endpoint.",
        x2 + 0.38, row_y + 5.65, col_w - 0.76, 1.1, SZ_BODY, False, SLATE,
    )
    add_text(slide, "ŷ = f_base(C) + Δ(C)", x2 + 0.55, row_y + 6.9, col_w - 1.1, 0.7, SZ_SUB, True, GREEN, PP_ALIGN.CENTER)

    add_panel(slide, "6  EXPERIMENTAL COMPARISON", x3, row_y, col_w, 7.95)
    add_progression_chart(slide, x3 + 0.6, row_y + 1.05, col_w - 1.2, 2.8)
    add_text(slide, "Scale + timbre + residual progression (corr)", x3 + 0.35, row_y + 4.2, col_w - 0.7, 0.55, SZ_BODY, True, MUTED, PP_ALIGN.CENTER)
    add_text(
        slide,
        "10k visual prior: 0.507. Final 59k prompt-stats residual: 0.582. Timbre/audio-prompt conditioning and residual recon produce the larger gains.",
        x3 + 0.4, row_y + 5.0, col_w - 0.8, 1.5, SZ_BODY, False, SLATE,
    )
    add_tag(slide, "final corr 0.5818 | MSE 0.6676 | MAE 0.6018", x3 + 0.65, row_y + 7.0, col_w - 1.3, GREEN)

    # Raw-video inference and Trump visual demo.
    demo_y = 27.85
    add_panel(slide, "7  RAW / SILENT VIDEO PIPELINE", 1.05, demo_y, 12.2, 7.5)
    add_fit_image(slide, ASSET_DIR / "report_silent_ref_pipeline.png", 1.45, demo_y + 0.75, 11.35, 6.65)

    trump_x = 13.65
    trump_w = 18.4
    add_panel(slide, "8  TRUMP SILENT → VOICED DEMO", trump_x, demo_y, trump_w, 7.5)
    seq = ASSET_DIR / "trump_sequence"
    add_frame_strip(slide, "silent input sequence",
        [seq / "silent_1.png", seq / "silent_2.png", seq / "silent_3.png"],
        trump_x + 0.55, demo_y + 0.8, 5.0, 2.2, TEAL)
    add_text(slide, "+", trump_x + 5.75, demo_y + 1.7, 0.4, 0.45, SZ_SUB, True, TEAL, PP_ALIGN.CENTER)
    add_frame_strip(slide, "reference tail",
        [seq / "ref_1.png", seq / "ref_2.png", seq / "ref_3.png"],
        trump_x + 6.25, demo_y + 0.8, 5.0, 2.2, ORANGE)
    add_text(slide, "→", trump_x + 11.45, demo_y + 1.7, 0.5, 0.45, SZ_SUB, True, TEAL, PP_ALIGN.CENTER)
    add_frame_strip(slide, "generated output",
        [seq / "gen_1.png", seq / "gen_2.png", seq / "gen_3.png"],
        trump_x + 12.1, demo_y + 0.8, 5.0, 2.2, GREEN)
    add_fit_image(slide, ASSET_DIR / "trump_ref_wave.png", trump_x + 0.9, demo_y + 3.3, 7.1, 1.2)
    add_text(slide, "reference audio prompt", trump_x + 0.9, demo_y + 4.55, 7.1, 0.5, SZ_BODY, True, MUTED, PP_ALIGN.CENTER)
    add_fit_image(slide, ASSET_DIR / "trump_generated_wave.png", trump_x + 9.3, demo_y + 3.3, 7.85, 1.2)
    add_text(slide, "generated waveform", trump_x + 9.3, demo_y + 4.55, 7.85, 0.5, SZ_BODY, True, MUTED, PP_ALIGN.CENTER)
    add_tag(slide, "silent video + reference audio → AVSR text → Mimi recon → voiced MP4", trump_x + 0.75, demo_y + 5.3, trump_w - 1.5, GREEN)
    add_text(slide, "Trump corr: 0.3035 → 0.3869 with final lip-AVSR preprocessing path.",
             trump_x + 0.65, demo_y + 6.25, trump_w - 1.3, 0.6, SZ_BODY, True, NAVY, PP_ALIGN.CENTER)

    # StreamLip V5 VSR panel — full width
    v5_y = 35.85
    add_panel(slide, "9  STREAMLIP V5: SELF-TRAINED VSR BRANCH", 1.05, v5_y, 31.0, 5.0)
    add_text(slide, "OLMo-1B + frozen Conformer + Gated Cross-Attention  (LRS3 pretrain split)",
             1.45, v5_y + 0.72, 14.0, 0.6, SZ_BODY, True, TEAL)
    v5_rows = [
        ("System", "WER", "Word Acc", "Notes"),
        ("External VSR baseline (beam=40)", "20.2%", "79.8%", "CTC+Att, 5k vocab"),
        ("StreamLip V5 (beam=3, step 4500)", "29.2%", "70.8%", "OLMo-1B, 50k vocab"),
    ]
    add_mini_table(slide, v5_rows[0], v5_rows[1:], 1.45, v5_y + 1.45, 14.6, row_h=0.55, size=SZ_SMALL)
    add_text(slide, "Where V5 beats the baseline:", 16.5, v5_y + 0.72, 14.6, 0.6, SZ_BODY, True, GREEN)
    add_bullets(
        slide,
        [
            "▸ Long sentences (≥10w, ≥4s): V5 win rate 10-16% vs 6-8% overall",
            "▸ LM prior corrects low-freq words, collocations & contractions (we're, they've)",
            "▸ Beam=3 optimal — wider beam hallucinates visually-ungrounded candidates",
        ],
        16.5, v5_y + 1.45, 14.6, 2.8, SZ_BODY, SLATE,
    )
    add_tag(slide, "V5 vs baseline on audio recon: Δcorr = 0.0007  |  alignment dominates transcript accuracy",
            2.2, v5_y + 4.45, 28.7, TEAL)

    # Bottom takeaway row
    take_y = v5_y + 5.0 + 0.33
    take_h = 41.18 - take_y - 0.1
    add_panel(slide, "10  TAKEAWAY", 1.05, take_y, 15.15, take_h)
    add_text(
        slide,
        "Do not cascade lip → text → TTS.",
        1.45, take_y + 0.72, 14.35, 0.9, SZ_SUB, True, TEAL,
    )
    add_bullets(
        slide,
        [
            "▸ Lip motion provides phonetic and timing evidence.",
            "▸ Mimi latents preserve codec-compatible acoustic structure.",
            "▸ Weak text stabilizes content — does not need to be perfect.",
            "▸ Timbre/reference conditions carry speaker and recording color.",
            "▸ StreamLip V5: LM-based VSR can serve as self-trained text branch.",
        ],
        1.45, take_y + 1.72, 14.35, take_h - 1.9, SZ_BODY, SLATE,
    )

    add_panel(slide, "11  REPRODUCE / ARTIFACTS", 16.9, take_y, 15.15, take_h)
    add_text(slide, "One-command raw-video demo:", 17.3, take_y + 0.72, 14.4, 0.6, SZ_BODY, True, TEAL)
    add_text(
        slide,
        "python scripts/run_raw_video_avsr_recon_pipeline.py\n  --input data/trump.mov --exp demo --force",
        17.3, take_y + 1.42, 14.35, 1.1, SZ_SMALL, False, NAVY,
    )
    card_w = 6.95
    add_artifact_card(slide, "Gradio GUI", "scripts/gradio_avsr_gui.py", 17.3, take_y + 2.65, card_w, 1.0, TEAL)
    add_artifact_card(slide, "checkpoint", "timbre-fix residual recon  step_002000.pt", 24.72, take_y + 2.65, card_w, 1.0, GREEN)

    prs.save(str(OUT))
    print(OUT)


if __name__ == "__main__":
    build()
