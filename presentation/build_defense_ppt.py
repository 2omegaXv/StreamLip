#!/usr/bin/env python3
"""Build the StreamLip defense presentation from report/poster assets."""

from __future__ import annotations

import hashlib
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "presentation"
REPORT_ASSETS = ROOT / "report" / "assets"
POSTER_ASSETS = ROOT / "poster" / "assets"
OUT = OUT_DIR / "streamlip_defense.pptx"
NOTES = OUT_DIR / "streamlip_defense_speaker_notes_zh.md"
FORMULA_DIR = OUT_DIR / "formula_assets"


NAVY = RGBColor(23, 35, 57)
PURPLE = RGBColor(90, 35, 112)
PURPLE2 = RGBColor(129, 66, 139)
PLUM = RGBColor(121, 68, 121)
TEAL = RGBColor(38, 104, 119)
GREEN = RGBColor(72, 126, 111)
ORANGE = RGBColor(158, 118, 54)
RED = RGBColor(137, 65, 91)
SLATE = RGBColor(49, 61, 80)
MUTED = RGBColor(101, 113, 132)
BORDER = RGBColor(205, 213, 224)
LIGHT = RGBColor(246, 248, 251)
WHITE = RGBColor(255, 255, 255)
PANEL_FILL = RGBColor(251, 252, 254)
HEADER_FILL = RGBColor(238, 241, 247)


def inch(v: float):
    return Inches(v)


def add_bg(slide):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = LIGHT
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, inch(0), inch(0), inch(13.333), inch(0.28))
    band.fill.solid()
    band.fill.fore_color.rgb = PURPLE
    band.line.fill.background()
    accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, inch(0), inch(0.28), inch(13.333), inch(0.04))
    accent.fill.solid()
    accent.fill.fore_color.rgb = PURPLE2
    accent.line.fill.background()


def set_text(tf, text, size=20, bold=False, color=SLATE, align=None):
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Arial"
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    if align is not None:
        p.alignment = align
    return p


def add_text(slide, text, x, y, w, h, size=20, bold=False, color=SLATE, align=None):
    box = slide.shapes.add_textbox(inch(x), inch(y), inch(w), inch(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = inch(0.03)
    tf.margin_right = inch(0.03)
    tf.margin_top = inch(0.03)
    tf.margin_bottom = inch(0.03)
    set_text(tf, text, size, bold, color, align)
    return box


def add_title(slide, title, subtitle=None):
    add_text(slide, title, 0.58, 0.55, 12.2, 0.55, 26, True, NAVY)
    if subtitle:
        add_text(slide, subtitle, 0.6, 1.08, 12.0, 0.34, 12, False, MUTED)


def add_footer(slide, idx):
    add_text(slide, "StreamLip Defense Presentation", 0.6, 7.14, 4.2, 0.18, 8, False, MUTED)
    add_text(slide, str(idx), 12.1, 7.14, 0.6, 0.18, 8, False, MUTED, PP_ALIGN.RIGHT)


def add_panel(slide, title, x, y, w, h, accent=TEAL):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, inch(x), inch(y), inch(w), inch(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PANEL_FILL
    shape.line.color.rgb = BORDER
    shape.line.width = Pt(1.0)
    shape.adjustments[0] = 0.04
    add_text(slide, title, x + 0.18, y + 0.12, w - 0.36, 0.25, 11, True, accent)
    return shape


def add_tag(slide, text, x, y, w, color=TEAL, size=10):
    tag = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, inch(x), inch(y), inch(w), inch(0.34))
    tag.fill.solid()
    tag.fill.fore_color.rgb = color
    tag.line.fill.background()
    tag.adjustments[0] = 0.12
    add_text(slide, text, x + 0.08, y + 0.08, w - 0.16, 0.16, size, True, WHITE, PP_ALIGN.CENTER)


def add_bullets(slide, items, x, y, w, h, size=15, color=SLATE):
    box = slide.shapes.add_textbox(inch(x), inch(y), inch(w), inch(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = inch(0.04)
    tf.margin_right = inch(0.04)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.name = "Arial"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(8)
    return box


def fit_image(slide, path, x, y, w, h):
    if not path.exists():
        add_panel(slide, f"Missing: {path.name}", x, y, w, h, RED)
        return
    with Image.open(path) as img:
        ratio = img.width / img.height
    box_ratio = w / h
    if ratio > box_ratio:
        dw = w
        dh = w / ratio
    else:
        dh = h
        dw = h * ratio
    dx = x + (w - dw) / 2
    dy = y + (h - dh) / 2
    slide.shapes.add_picture(str(path), inch(dx), inch(dy), inch(dw), inch(dh))


def formula_path(expr: str, fontsize=22, color="#172339") -> Path:
    FORMULA_DIR.mkdir(parents=True, exist_ok=True)
    key = hashlib.sha1(f"{expr}|{fontsize}|{color}".encode("utf-8")).hexdigest()[:16]
    path = FORMULA_DIR / f"formula_{key}.png"
    if path.exists():
        return path
    fig = plt.figure(figsize=(0.01, 0.01), dpi=240)
    text = fig.text(0, 0, expr, fontsize=fontsize, color=color, usetex=False)
    fig.canvas.draw()
    bbox = text.get_window_extent()
    width = max(1, bbox.width / fig.dpi)
    height = max(1, bbox.height / fig.dpi)
    plt.close(fig)

    fig = plt.figure(figsize=(width + 0.18, height + 0.12), dpi=240)
    fig.patch.set_alpha(0)
    fig.text(0.06 / (width + 0.18), 0.5, expr, fontsize=fontsize, color=color, va="center", usetex=False)
    fig.savefig(path, transparent=True, bbox_inches="tight", pad_inches=0.03)
    plt.close(fig)
    return path


def add_formula(slide, expr, x, y, w, h, fontsize=22):
    fit_image(slide, formula_path(expr, fontsize=fontsize), x, y, w, h)


def add_formula_panel(slide, title, expr, x, y, w, h, accent=PURPLE, fontsize=22):
    add_panel(slide, title, x, y, w, h, accent)
    formula_y = y + 0.46
    formula_h = max(0.16, h - 0.58)
    add_formula(slide, expr, x + 0.26, formula_y, w - 0.52, formula_h, fontsize=fontsize)


def add_node(slide, text, x, y, w, h, accent=TEAL, fill=WHITE, size=9):
    node = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, inch(x), inch(y), inch(w), inch(h))
    node.fill.solid()
    node.fill.fore_color.rgb = fill
    node.line.color.rgb = accent
    node.line.width = Pt(1.0)
    node.adjustments[0] = 0.08
    add_text(slide, text, x + 0.08, y + 0.08, w - 0.16, h - 0.14, size, True, NAVY, PP_ALIGN.CENTER)
    return node


def add_arrow_label(slide, x, y, color=MUTED):
    add_text(slide, ">", x, y, 0.2, 0.2, 12, True, color, PP_ALIGN.CENTER)


def add_recon_diagram(slide):
    add_panel(slide, "rejected stochastic sampling", 0.85, 1.35, 3.45, 4.15, PLUM)
    add_formula(slide, r"$x_0\sim\mathcal{N}(0,I)$", 1.08, 1.9, 1.15, 0.28, 14)
    add_formula(slide, r"$\tau:0\rightarrow1$", 2.65, 1.9, 1.0, 0.28, 14)
    for i, (cx, cy, color) in enumerate([(1.3, 2.75, PLUM), (2.05, 3.18, TEAL), (2.8, 3.62, PURPLE)]):
        dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, inch(cx), inch(cy), inch(0.2), inch(0.2))
        dot.fill.solid()
        dot.fill.fore_color.rgb = color
        dot.line.fill.background()
        if i < 2:
            add_arrow_label(slide, cx + 0.36, cy + 0.05, MUTED)
    add_text(slide, "iterative denoise trajectory", 1.1, 4.18, 2.8, 0.24, 11, True, MUTED, PP_ALIGN.CENTER)
    add_text(slide, "unnecessary sampling loop for paired reconstruction", 1.15, 4.62, 2.65, 0.34, 9, False, SLATE, PP_ALIGN.CENTER)

    add_panel(slide, "adopted endpoint reconstruction", 4.55, 1.35, 3.55, 4.15, PURPLE)
    add_node(slide, "condition set\nC", 4.85, 2.0, 0.8, 0.62, PURPLE, HEADER_FILL, 9)
    add_arrow_label(slide, 5.72, 2.19, MUTED)
    add_node(slide, "frozen base\nf_base(C)", 6.05, 1.82, 1.25, 0.55, PURPLE, WHITE, 8)
    add_node(slide, "residual\nf_res(C)", 6.05, 2.72, 1.25, 0.55, TEAL, WHITE, 8)
    add_formula(slide, r"$y_{\mathrm{base}}$", 7.15, 1.92, 0.48, 0.18, 10)
    add_formula(slide, r"$\Delta$", 7.25, 2.85, 0.32, 0.18, 10)
    add_text(slide, "+", 7.25, 2.35, 0.22, 0.2, 15, True, NAVY, PP_ALIGN.CENTER)
    add_arrow_label(slide, 7.48, 2.35, MUTED)
    add_node(slide, "output\nlatent", 7.68, 2.22, 0.62, 0.52, PURPLE, HEADER_FILL, 8)
    add_formula(slide, r"$\hat y$", 7.79, 2.85, 0.32, 0.18, 10)
    add_formula(slide, r"$\hat y=y_{\mathrm{base}}+\Delta$", 5.05, 3.65, 2.45, 0.28, 15)
    add_text(slide, "direct endpoint regression in normalized Mimi latent space", 4.92, 4.5, 2.8, 0.34, 9, False, SLATE, PP_ALIGN.CENTER)


def add_metric(slide, label, value, x, y, w, color=TEAL):
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, inch(x), inch(y), inch(w), inch(0.92))
    card.fill.solid()
    card.fill.fore_color.rgb = PANEL_FILL
    card.line.color.rgb = BORDER
    card.adjustments[0] = 0.08
    add_text(slide, value, x + 0.1, y + 0.13, w - 0.2, 0.32, 20, True, color, PP_ALIGN.CENTER)
    add_text(slide, label, x + 0.1, y + 0.52, w - 0.2, 0.2, 9, True, MUTED, PP_ALIGN.CENTER)


def add_flow(slide, steps, x, y, w, h):
    gap = 0.12
    step_w = (w - gap * (len(steps) - 1)) / len(steps)
    colors = [TEAL, PLUM, GREEN, NAVY, PURPLE2]
    for i, step in enumerate(steps):
        sx = x + i * (step_w + gap)
        node = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, inch(sx), inch(y), inch(step_w), inch(h))
        node.fill.solid()
        node.fill.fore_color.rgb = colors[i % len(colors)]
        node.line.fill.background()
        node.adjustments[0] = 0.08
        add_text(slide, step, sx + 0.08, y + 0.17, step_w - 0.16, h - 0.24, 13, True, WHITE, PP_ALIGN.CENTER)
        if i < len(steps) - 1:
            add_text(slide, ">", sx + step_w - 0.02, y + h / 2 - 0.11, 0.22, 0.2, 12, True, MUTED, PP_ALIGN.CENTER)


def add_table(slide, headers, rows, x, y, w, row_h=0.34, size=9, widths=None):
    n = len(headers)
    widths = widths or [w / n] * n
    rect = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, inch(x), inch(y), inch(w), inch(row_h))
    rect.fill.solid()
    rect.fill.fore_color.rgb = HEADER_FILL
    rect.line.color.rgb = BORDER
    cx = x
    for head, cw in zip(headers, widths):
        add_text(slide, head, cx + 0.05, y + 0.08, cw - 0.1, 0.14, size, True, NAVY)
        cx += cw
    for r, row in enumerate(rows):
        yy = y + row_h * (r + 1)
        fill = RGBColor(248, 250, 253) if r % 2 == 0 else WHITE
        rect = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, inch(x), inch(yy), inch(w), inch(row_h))
        rect.fill.solid()
        rect.fill.fore_color.rgb = fill
        rect.line.color.rgb = RGBColor(226, 232, 240)
        cx = x
        for cell, cw in zip(row, widths):
            add_text(slide, cell, cx + 0.05, yy + 0.08, cw - 0.1, 0.14, size, False, SLATE)
            cx += cw


def add_chart(slide, x, y, w, h):
    labels = ["10k", "30k", "+timbre", "+prompt/res", "59k final"]
    values = [0.507, 0.531, 0.563, 0.573, 0.582]
    max_v = 0.60
    gap = 0.14
    bw = (w - gap * (len(values) - 1)) / len(values)
    for i, (lab, val) in enumerate(zip(labels, values)):
        bx = x + i * (bw + gap)
        bh = h * val / max_v
        by = y + h - bh
        color = PURPLE if i == len(values) - 1 else (TEAL if i >= 2 else PLUM)
        bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, inch(bx), inch(by), inch(bw), inch(bh))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        add_text(slide, f"{val:.3f}", bx - 0.02, by - 0.26, bw + 0.04, 0.16, 8, True, NAVY, PP_ALIGN.CENTER)
        add_text(slide, lab, bx - 0.05, y + h + 0.07, bw + 0.1, 0.16, 8, False, MUTED, PP_ALIGN.CENTER)
    axis = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, inch(x), inch(y + h), inch(w), inch(0.02))
    axis.fill.solid()
    axis.fill.fore_color.rgb = BORDER
    axis.line.fill.background()


SLIDES = [
    {
        "title": "StreamLip: Audio-Centric Visual Speech Reconstruction",
        "note": "开场先不要讲模型细节，直接说明任务：我们想从被静音、被屏蔽或音频缺失的视频里恢复说话声音。核心不是做字幕，而是恢复和嘴型同步的语音。",
    },
    {
        "title": "Opening Problem",
        "note": "这一页把问题抛出来：如果访谈或对话视频的音频缺失，只有脸和嘴型，我们如何恢复与画面同步的语音？强调这是一个音频恢复问题，不是单纯识别文字。",
    },
    {
        "title": "Why Not Video -> Text -> TTS?",
        "note": "解释 cascade 的问题：文字会丢掉音色、韵律、发音细节，也很难保证 frame-level audio-video sync。识别错误还会在 TTS 阶段放大，所以我们需要直接处理视频到音频表示的关系。",
    },
    {
        "title": "Core Idea",
        "note": "给出本项目的主张：text 只是弱语义和对齐条件，主体路线是 lip/video 加 timbre 直接到 Mimi audio latent，再解码成 waveform。",
    },
    {
        "title": "Final System Overview",
        "note": "用架构图讲系统。左侧是视觉、文本、说话人/音色条件，右侧预测 normalized Mimi latent，然后 Mimi decoder 输出 waveform。重点说所有条件都对 recon head 可见。",
    },
    {
        "title": "Data and Target",
        "note": "讲数据不是只存文本，而是每个 clip 都构造成视觉、文本 hidden、speaker、Mimi target、prompt/timbre 的配套表示。训练集 59144，验证集固定 1000。",
    },
    {
        "title": "Audio-Latent Formulation",
        "note": "解释为什么用 Mimi latent。我们不是直接回归 waveform，也不是从 transcript 做 TTS，而是在 codec-compatible 的连续音频 latent 空间里训练和评估。",
    },
    {
        "title": "Residual Endpoint Reconstruction",
        "note": "讲最终放弃随机 flow matching sampling。固定 x_tilde=0、tau=1，直接做 endpoint reconstruction，并用 frozen base 加 residual correction 提升结果。训练目标由三部分组成：latent MSE 负责数值重构，相关性项鼓励整体时序形状一致，prompt statistics 项约束输出的均值/方差接近参考音频的音色统计。",
    },
    {
        "title": "Self-Trained Visual Text Branch",
        "note": "这一页讲 StreamLip V5，但不要把它变成项目中心。它证明文本分支可以内部训练。指标是 WER 29.2、word accuracy 70.8，训练得可以，但最终只是弱 text conditioning。",
    },
    {
        "title": "Text Is Necessary, But Need Not Be Perfect",
        "note": "这里讲最关键实验：完全去掉 text condition 会明显下降，所以 text 必须有；但 text 不需要非常准，GT text 到 decoded text 只掉 0.0035 corr。即使文字错误较高，人耳仍可能听出大概，因为 lip/audio latent 条件还在。",
    },
    {
        "title": "Main Experimental Progression",
        "note": "用柱状图讲从 10k visual prior 到最终 59k prompt-stats residual 的进步。重点是 data scale、timbre/audio prompt、residual recon 是主要收益来源。",
    },
    {
        "title": "Trump Demo Case",
        "note": "这里只讲一个 demo case：输入 silent Trump video，给一段 reference audio，输出合成后的 voiced video。不要展开 pipeline 细节，用它证明系统可以用于缺失音频的视频恢复。",
    },
    {
        "title": "Takeaway",
        "note": "最后收束：恢复缺失音频不能只靠 text cascade，必须直接重建 audio latent；text 有用但应是弱条件；同步、音色和声学细节需要 lip/audio/timbre 条件保留在 reconstructor 中。",
    },
]


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    notes = ["# StreamLip Defense Speaker Notes\n"]

    # Slide 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_text(slide, "StreamLip", 0.78, 1.05, 12.0, 0.62, 38, True, PURPLE, PP_ALIGN.CENTER)
    add_text(slide, "Audio-Centric Visual Speech Reconstruction", 1.2, 1.82, 11.0, 0.45, 24, True, NAVY, PP_ALIGN.CENTER)
    add_text(slide, "with Weak Text Conditioning", 1.2, 2.36, 11.0, 0.35, 18, True, TEAL, PP_ALIGN.CENTER)
    add_flow(slide, ["silent / corrupted video", "Mimi audio latents", "recovered speech"], 1.6, 3.35, 10.2, 0.78)
    add_text(slide, "DL-V2A Project Team | fm-avsr-cleanup", 1.2, 5.15, 11.0, 0.25, 13, False, MUTED, PP_ALIGN.CENTER)
    add_footer(slide, 1)
    notes.append("## 1. Title\n" + SLIDES[0]["note"] + "\n")

    # Slide 2
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide); add_title(slide, "Opening Problem", "Recover speech when audio is muted, corrupted, or missing")
    add_text(slide, "How can we recover speech from an interview/dialogue video when the audio track is unavailable?", 0.8, 1.55, 11.8, 0.55, 24, True, NAVY, PP_ALIGN.CENTER)
    add_panel(slide, "Given", 0.85, 2.55, 3.6, 2.4, TEAL)
    add_bullets(slide, ["Talking-face video", "Lip motion and face identity", "No usable speech track"], 1.15, 3.1, 3.0, 1.3, 15)
    add_panel(slide, "Need", 4.85, 2.55, 3.6, 2.4, PURPLE)
    add_bullets(slide, ["Speech waveform", "Aligned to mouth motion", "Speaker/timbre consistency"], 5.15, 3.1, 3.0, 1.3, 15)
    add_panel(slide, "Challenge", 8.85, 2.55, 3.6, 2.4, PLUM)
    add_bullets(slide, ["Text is incomplete", "Audio detail is not textual", "Synchronization is frame-level"], 9.15, 3.1, 3.0, 1.3, 15)
    add_tag(slide, "This is audio recovery, not only speech recognition.", 3.0, 5.55, 7.4, TEAL, 12)
    add_footer(slide, 2)
    notes.append("## 2. Opening Problem\n" + SLIDES[1]["note"] + "\n")

    # Slide 3
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide); add_title(slide, "Why Not Video -> Text -> TTS?", "A cascade discards information the audio needs")
    add_flow(slide, ["video", "text", "TTS", "speech"], 1.0, 1.65, 11.2, 0.68)
    add_panel(slide, "What text loses", 0.85, 2.75, 3.8, 2.4, PLUM)
    add_bullets(slide, ["phonetic timing", "prosody and rhythm", "speaker color", "non-text acoustic cues"], 1.15, 3.22, 3.2, 1.4, 14)
    add_panel(slide, "What TTS cannot guarantee", 4.8, 2.75, 3.8, 2.4, PURPLE)
    add_bullets(slide, ["frame-level lip/audio sync", "exact duration and pacing", "robustness to recognition errors"], 5.1, 3.22, 3.2, 1.4, 14)
    add_panel(slide, "Our response", 8.75, 2.75, 3.8, 2.4, TEAL)
    add_bullets(slide, ["keep lip evidence visible", "condition on timbre directly", "predict audio latents directly"], 9.05, 3.22, 3.2, 1.4, 14)
    add_tag(slide, "Directly model the video/audio-latent relation.", 3.3, 5.68, 6.7, PURPLE, 12)
    add_footer(slide, 3)
    notes.append("## 3. Why Not Video -> Text -> TTS?\n" + SLIDES[2]["note"] + "\n")

    # Slide 4
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide); add_title(slide, "Core Idea", "Audio-centric reconstruction with weak text conditioning")
    add_formula_panel(
        slide,
        "conditioning set and predictor",
        r"$C=\{e^{v}_{1:T_v},h_{1:L},s,A,q\},\qquad \hat y_{1:T}=f_{\theta}(C),\qquad \hat x_{\mathrm{wav}}=\mathrm{MimiDecoder}(\mathrm{denorm}(\hat y))$",
        0.9,
        1.55,
        11.55,
        1.35,
        PURPLE,
        15,
    )
    add_metric(slide, "Final corr", "0.5818", 1.4, 3.55, 2.2, PURPLE)
    add_metric(slide, "GT -> decoded text drop", "0.0035", 4.1, 3.55, 2.8, TEAL)
    add_metric(slide, "Train clips", "59,144", 7.4, 3.55, 2.2, PLUM)
    add_metric(slide, "Val clips", "1,000", 10.1, 3.55, 2.0, GREEN)
    add_text(slide, "Text helps content and alignment, but the model is not forced to synthesize speech from text alone.", 1.2, 5.25, 10.8, 0.38, 18, True, NAVY, PP_ALIGN.CENTER)
    add_footer(slide, 4)
    notes.append("## 4. Core Idea\n" + SLIDES[3]["note"] + "\n")

    # Slide 5
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide); add_title(slide, "Final System Overview", "All key conditions remain visible to the reconstructor")
    fit_image(slide, REPORT_ASSETS / "system_architecture.png", 0.75, 1.35, 11.85, 5.25)
    add_footer(slide, 5)
    notes.append("## 5. Final System Overview\n" + SLIDES[4]["note"] + "\n")

    # Slide 6
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide); add_title(slide, "Data and Target", "Each clip is converted into matched visual, text, timbre, and audio-latent files")
    add_metric(slide, "training clips", "59,144", 0.85, 1.35, 2.3, TEAL)
    add_metric(slide, "validation clips", "1,000", 3.45, 1.35, 2.3, GREEN)
    add_metric(slide, "prompt length", "38 frames", 6.05, 1.35, 2.3, ORANGE)
    add_metric(slide, "prompt duration", "3.04 s", 8.65, 1.35, 2.3, PURPLE)
    headers = ["Representation", "Role"]
    rows = [
        ["lip_avsr.npy", "StreamLip-compatible mouth crop"],
        ["SmolLM2 hidden", "weak semantic/alignment condition"],
        ["speaker_emb.npy", "face-based speaker identity"],
        ["audio_prompt.npy", "Mimi prompt tokens"],
        ["timbre_cond.npy", "prompt mean/std timbre"],
        ["latent.npz", "Mimi latent training target"],
    ]
    add_table(slide, headers, rows, 1.0, 2.7, 11.2, row_h=0.42, size=10, widths=[3.0, 8.2])
    add_tag(slide, "All metrics are computed in normalized Mimi latent space.", 3.25, 5.75, 6.8, TEAL, 11)
    add_footer(slide, 6)
    notes.append("## 6. Data and Target\n" + SLIDES[5]["note"] + "\n")

    # Slide 7
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide); add_title(slide, "Audio-Latent Formulation", "Recover codec-compatible acoustic structure rather than only words")
    add_formula_panel(slide, "normalized Mimi target", r"$z_t\in\mathbb{R}^{512},\qquad y_t=\frac{z_t-\mu}{\sigma}$", 0.9, 1.35, 5.4, 2.0, PURPLE, 24)
    add_formula_panel(slide, "denormalize and decode", r"$\hat z_t=\hat y_t\sigma+\mu,\qquad \hat x_{\mathrm{wav}}=\mathrm{MimiDecoder}(\hat z_{1:T})$", 7.0, 1.35, 5.4, 2.0, TEAL, 19)
    add_bullets(slide, [
        "Latent prediction preserves acoustic information beyond text.",
        "The objective directly supervises paired audio reconstruction.",
        "This helps maintain timing and speaker/timbre cues visible in video/audio prompts.",
    ], 1.25, 4.0, 10.8, 1.45, 17)
    add_tag(slide, "Not waveform regression. Not transcript-to-TTS.", 3.55, 5.88, 6.2, PURPLE, 12)
    add_footer(slide, 7)
    notes.append("## 7. Audio-Latent Formulation\n" + SLIDES[6]["note"] + "\n")

    # Slide 8
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide); add_title(slide, "Residual Endpoint Reconstruction", "The final route is deterministic reconstruction, not stochastic sampling")
    add_recon_diagram(slide)
    add_formula_panel(slide, "deterministic endpoint", r"$\tilde{x}=0,\quad \tau=1$", 8.55, 1.25, 3.85, 0.86, PURPLE, 24)
    add_formula_panel(slide, "residual predictor", r"$y_{\mathrm{base}}=f_{\mathrm{base}}(C),\quad \Delta=f_{\mathrm{res}}(C),\quad \hat y=y_{\mathrm{base}}+\Delta$", 8.55, 2.23, 3.85, 0.9, TEAL, 13)
    add_formula_panel(slide, "training objective", r"$\mathcal{L}=\mathcal{L}_{\mathrm{rec}}+0.2\,\mathcal{L}_{\mathrm{corr}}+0.05\,\mathcal{L}_{\mathrm{stat}}$", 8.55, 3.25, 3.85, 0.86, PLUM, 16)
    add_panel(slide, "loss terms", 8.55, 4.25, 3.85, 1.65, GREEN)
    add_formula(slide, r"$\mathcal{L}_{\mathrm{rec}}=\frac{1}{T}\sum_{t=1}^{T}\|\hat y_t-y_t\|_2^2$", 8.88, 4.78, 3.15, 0.3, 12)
    add_formula(slide, r"$\mathcal{L}_{\mathrm{corr}}=1-\mathrm{corr}\!\left(\mathrm{vec}(\hat y_{39:T}),\mathrm{vec}(y_{39:T})\right)$", 8.82, 5.23, 3.25, 0.28, 10)
    add_formula(slide, r"$\mathcal{L}_{\mathrm{stat}}=\|\mu(\hat y_{1:38})-\mu(p)\|_2^2+\|\sigma(\hat y_{1:38})-\sigma(p)\|_2^2$", 8.82, 5.62, 3.25, 0.25, 8)
    add_footer(slide, 8)
    notes.append("## 8. Residual Endpoint Reconstruction\n" + SLIDES[7]["note"] + "\n")

    # Slide 9
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide); add_title(slide, "Self-Trained Visual Text Branch", "A useful internal weak-text source, not the center of the audio model")
    add_flow(slide, ["lip frames", "frozen visual encoder", "gated cross-attn", "OLMo-1B text"], 0.9, 1.55, 11.5, 0.72)
    add_metric(slide, "WER", "29.2%", 1.4, 3.0, 2.2, PLUM)
    add_metric(slide, "Word accuracy", "70.8%", 4.0, 3.0, 2.4, PURPLE)
    add_metric(slide, "Beam", "3", 6.9, 3.0, 1.7, TEAL)
    add_metric(slide, "Model size", "1.24B", 9.1, 3.0, 2.2, PURPLE)
    add_bullets(slide, [
        "Frozen visual features are injected into OLMo-1B via gated cross-attention.",
        "Large-LM context helps longer sentences; short clips still suffer insertion errors.",
        "The audio reconstructor remains robust because text is weak conditioning.",
    ], 1.15, 4.45, 11.2, 1.25, 15)
    add_footer(slide, 9)
    notes.append("## 9. Self-Trained Visual Text Branch\n" + SLIDES[8]["note"] + "\n")

    # Slide 10
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide); add_title(slide, "Text Is Necessary, But Need Not Be Perfect", "Text carries semantics/alignment, not all audio detail")
    headers = ["Text condition", "Alignment", "Corr", "MSE", "MAE"]
    rows = [
        ["GT text", "word ts", "0.58184", "0.66763", "0.60182"],
        ["Decoded text", "word ts", "0.57833", "0.67162", "0.60371"],
        ["Decoded text", "uniform", "0.28679", "1.68438", "1.02352"],
        ["StreamLip text", "uniform", "0.28609", "1.68488", "1.02375"],
    ]
    add_table(slide, headers, rows, 0.85, 1.45, 11.65, row_h=0.46, size=10, widths=[3.1, 2.1, 2.0, 2.0, 2.0])
    add_metric(slide, "GT -> decoded corr drop", "0.0035", 1.25, 4.25, 2.85, TEAL)
    add_metric(slide, "word-ts -> uniform drop", "~0.29", 4.55, 4.25, 2.85, PLUM)
    add_metric(slide, "Text condition", "required", 7.85, 4.25, 2.85, PURPLE)
    add_formula(slide, r"$\Delta_{\mathrm{text}}=0.58184-0.57833=0.00351$", 1.2, 5.18, 4.2, 0.38, 18)
    add_text(slide, "Ablations show text cannot simply be removed, but transcript accuracy is not the main bottleneck once lip/audio-latent conditions are available.", 1.1, 5.68, 11.1, 0.4, 16, True, NAVY, PP_ALIGN.CENTER)
    add_footer(slide, 10)
    notes.append("## 10. Text Is Necessary, But Need Not Be Perfect\n" + SLIDES[9]["note"] + "\n")

    # Slide 11
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide); add_title(slide, "Main Experimental Progression", "Scale helps; timbre/audio prompt and residual reconstruction help more")
    add_chart(slide, 0.95, 1.55, 7.1, 3.0)
    add_panel(slide, "Final checkpoint", 8.55, 1.45, 3.8, 2.05, PURPLE)
    add_formula(slide, r"$\mathrm{corr}=0.58184180$", 8.95, 1.9, 3.0, 0.32, 17)
    add_formula(slide, r"$\mathrm{MSE}=0.66763106$", 8.95, 2.33, 3.0, 0.32, 17)
    add_formula(slide, r"$\mathrm{MAE}=0.60181897$", 8.95, 2.76, 3.0, 0.32, 17)
    add_panel(slide, "Interpretation", 8.55, 3.85, 3.8, 1.85, TEAL)
    add_bullets(slide, ["data scale improves visual prior", "timbre gives larger gains", "residual recon stabilizes final route"], 8.85, 4.28, 3.2, 0.9, 12)
    add_footer(slide, 11)
    notes.append("## 11. Main Experimental Progression\n" + SLIDES[10]["note"] + "\n")

    # Slide 12
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide); add_title(slide, "Trump Demo Case", "Reference audio + silent video -> synthesized voiced video")
    seq = POSTER_ASSETS / "trump_sequence"
    add_panel(slide, "Reference audio", 0.75, 1.35, 3.85, 4.55, PLUM)
    fit_image(slide, seq / "ref_2.png", 1.55, 1.9, 2.2, 2.2)
    fit_image(slide, POSTER_ASSETS / "trump_ref_wave.png", 1.1, 4.45, 3.15, 0.62)
    add_panel(slide, "Silent video", 4.85, 1.35, 3.85, 4.55, TEAL)
    fit_image(slide, seq / "silent_2.png", 5.65, 2.05, 2.2, 2.2)
    add_text(slide, "mouth motion without usable speech", 5.1, 4.6, 3.35, 0.26, 12, True, MUTED, PP_ALIGN.CENTER)
    add_panel(slide, "Synthesized video", 8.95, 1.35, 3.85, 4.55, PURPLE)
    fit_image(slide, seq / "gen_2.png", 9.75, 2.05, 2.2, 2.2)
    fit_image(slide, POSTER_ASSETS / "trump_generated_wave.png", 9.3, 4.45, 3.15, 0.62)
    add_tag(slide, "case study: recover voiced output from silent visual input", 2.6, 6.25, 8.2, PURPLE, 12)
    add_footer(slide, 12)
    notes.append("## 12. Trump Demo Case\n" + SLIDES[11]["note"] + "\n")

    # Slide 13
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide); add_title(slide, "Takeaway", "The architectural conclusion")
    add_panel(slide, "1", 0.9, 1.55, 3.7, 3.25, TEAL)
    add_text(slide, "Recovering speech from missing-audio video requires direct audio-latent reconstruction.", 1.2, 2.35, 3.1, 0.9, 18, True, NAVY, PP_ALIGN.CENTER)
    add_panel(slide, "2", 4.85, 1.55, 3.7, 3.25, PURPLE)
    add_text(slide, "Text is useful, but it should remain weak conditioning rather than the bottleneck.", 5.15, 2.35, 3.1, 0.9, 18, True, NAVY, PP_ALIGN.CENTER)
    add_panel(slide, "3", 8.8, 1.55, 3.7, 3.25, PLUM)
    add_text(slide, "Lip/audio synchronization and timbre require visual/audio conditions to stay visible.", 9.1, 2.35, 3.1, 0.9, 18, True, NAVY, PP_ALIGN.CENTER)
    add_text(slide, "StreamLip is best understood as audio-centric visual speech reconstruction.", 1.0, 5.62, 11.2, 0.35, 20, True, PURPLE, PP_ALIGN.CENTER)
    add_footer(slide, 13)
    notes.append("## 13. Takeaway\n" + SLIDES[12]["note"] + "\n")

    prs.save(OUT)
    NOTES.write_text("\n".join(notes), encoding="utf-8")
    print(OUT)
    print(NOTES)


if __name__ == "__main__":
    build()
